import faiss
import os
import pickle
import logging
from typing import List
from pathlib import Path
import tempfile
import shutil
import zipfile
import tarfile

from pdfminer.high_level import extract_text as extract_pdf_text
from PIL import Image
import pytesseract
from pptx import Presentation
import docx2txt
import pandas as pd
import numpy as np

logger = logging.getLogger(__name__)

class Retriever:
    def __init__(self, embedder, documents_path='documents/', chunk_size=500, chunk_overlap=100):
        if not hasattr(embedder, 'embed') or not callable(embedder.embed):
            raise AttributeError("The provided embedder object must have an 'embed' method.")
        self.embedder = embedder
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.index = None
        self.documents_path = documents_path
        self.texts = []

    def _chunk_text(self, text: str) -> List[str]:
        return [text[i:i+self.chunk_size] for i in range(0, len(text), self.chunk_size - self.chunk_overlap)]

    def _extract_text(self, file_path: Path) -> str:
        ext = file_path.suffix.lower()
        if ext == '.pdf':
            return extract_pdf_text(str(file_path))
        elif ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
            image = Image.open(file_path)
            return pytesseract.image_to_string(image)
        elif ext == '.pptx':
            prs = Presentation(file_path)
            return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and isinstance(shape.text, str)])
        elif ext in ['.doc', '.docx']:
            return docx2txt.process(str(file_path))
        elif ext in ['.xls', '.xlsx']:
            try:
                df = pd.read_excel(file_path)
                return df.to_string(index=False)
            except Exception as e:
                logger.error(f"Error reading spreadsheet {file_path}: {e}")
                return ""
        elif ext in ['.txt', '.md', '.rst']:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        else:
            logger.warning(f"Unsupported file type: {file_path}")
            return ""

    def _extract_from_archive(self, archive_path: Path, temp_dir: Path) -> List[Path]:
        extracted_files = []
        try:
            if archive_path.suffix == '.zip':
                with zipfile.ZipFile(archive_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)
            elif archive_path.suffixes[-2:] == ['.tar', '.gz'] or archive_path.suffix == '.tar':
                with tarfile.open(archive_path, 'r:*') as tar:
                    tar.extractall(temp_dir)
            else:
                logger.warning(f"Unsupported archive type: {archive_path}")
                return []
            extracted_files = list(temp_dir.rglob("*"))
        except Exception as e:
            logger.error(f"Error extracting archive {archive_path}: {e}")
        return [f for f in extracted_files if f.is_file()]

    def index_documents(self, save_path: str = "index_store")-> bool:
        all_chunks = []
        folder_path = Path(self.documents_path)

        for file_path in folder_path.rglob("*"):
            if file_path.is_file():
                if file_path.suffix in ['.zip', '.tar', '.gz', '.tar.gz']:
                    with tempfile.TemporaryDirectory() as tmpdirname:
                        extracted_files = self._extract_from_archive(file_path, Path(tmpdirname))
                        for extracted_file in extracted_files:
                            logger.info(f"Processing extracted file: {extracted_file}")
                            text = self._extract_text(extracted_file)
                            if text:
                                chunks = self._chunk_text(text)
                                all_chunks.extend(chunks)
                else:
                    logger.info(f"Processing file: {file_path}")
                    text = self._extract_text(file_path)
                    if text:
                        chunks = self._chunk_text(text)
                        all_chunks.extend(chunks)

        embeddings = self.embedder.embed(all_chunks)
        self.index = faiss.IndexFlatL2(len(embeddings[0]))
        import numpy as np
        self.index.add(np.array(embeddings))
        self.texts = all_chunks

        os.makedirs(save_path, exist_ok=True)
        faiss.write_index(self.index, os.path.join(save_path, "faiss.index"))
        with open(os.path.join(save_path, "texts.pkl"), "wb") as f:
            pickle.dump(self.texts, f)
        return True

    def load_index(self, save_path: str = "index_store") -> bool:
        index_path = os.path.join(save_path, "faiss.index")
        texts_path = os.path.join(save_path, "texts.pkl")
        if not os.path.exists(index_path) or not os.path.exists(texts_path):
            raise FileNotFoundError("FAISS index or text storage file not found.")

        self.index = faiss.read_index(index_path)
        with open(texts_path, "rb") as f:
            self.texts = pickle.load(f)
        return True

    def retrieve(self, query: str, top_k: int = 5) -> List[str]:
        if self.index is None or not self.texts:
            raise ValueError("Index not initialized. Call 'index_documents' or 'load_index' first.")
        embedding = self.embedder.embed([query])[0]
        embedding_np = np.array([embedding], dtype='float32')  # Ensure shape (1, dim) and correct dtype
        _, I = self.index.search(embedding_np, top_k)
        retrieved = [self.texts[i] for i in I[0]]

        return retrieved
